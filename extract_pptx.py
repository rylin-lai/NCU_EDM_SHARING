#!/usr/bin/env python3

from pptx import Presentation
import sys

def extract_pptx_to_markdown(pptx_file, output_file):
    """Extract content from PPTX and convert to markdown outline."""
    prs = Presentation(pptx_file)
    
    markdown_content = []
    markdown_content.append("# Python-The-Effective-Tool-for-the-Workplace\n")
    
    for i, slide in enumerate(prs.slides, 1):
        slide_title = f"## Slide {i}"
        slide_content = []
        
        # Extract text from all shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    # If this looks like a title (usually the first non-empty text)
                    if not slide_content and len(text) < 100:
                        slide_title = f"## Slide {i}: {text}"
                    else:
                        slide_content.append(text)
        
        markdown_content.append(slide_title)
        markdown_content.append("")
        
        if slide_content:
            for content in slide_content:
                # Format as bullet points if content looks like a list item
                if content and not content.startswith("- ") and not content.startswith("* "):
                    # Split by newlines and format each as bullet point
                    lines = content.split('\n')
                    for line in lines:
                        line = line.strip()
                        if line:
                            markdown_content.append(f"- {line}")
                else:
                    markdown_content.append(content)
            markdown_content.append("")
        else:
            markdown_content.append("*(No text content)*")
            markdown_content.append("")
    
    # Write to markdown file
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write('\n'.join(markdown_content))
    
    print(f"Content extracted to {output_file}")

if __name__ == "__main__":
    pptx_file = "Python-The-Effective-Tool-for-the-Workplace.pptx"
    output_file = "大綱.md"
    
    try:
        extract_pptx_to_markdown(pptx_file, output_file)
    except Exception as e:
        print(f"Error: {e}")